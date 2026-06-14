"""
Script auxiliar (uso unico) para gerar reports/apresentacao.pptx a partir do
conteudo de reports/apresentacao.md, removendo as secoes de narracao/figuras
sugeridas/orcamento de tempo, que servem apenas como roteiro de apoio.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

COR_TITULO = RGBColor(0x1F, 0x3B, 0x4D)
COR_TEXTO = RGBColor(0x33, 0x33, 0x33)
COR_DESTAQUE = RGBColor(0x1F, 0x7A, 0x4D)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

LAYOUT_TITULO = prs.slide_layouts[0]
LAYOUT_CONTEUDO = prs.slide_layouts[1]
LAYOUT_SECAO = prs.slide_layouts[2]


def add_titulo_slide(titulo, subtitulos):
    slide = prs.slides.add_slide(LAYOUT_TITULO)
    slide.shapes.title.text = titulo
    slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(40)
    slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = COR_TITULO

    subtitle_ph = slide.placeholders[1]
    tf = subtitle_ph.text_frame
    tf.clear()
    for i, linha in enumerate(subtitulos):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = linha
        p.font.size = Pt(22)
        p.font.color.rgb = COR_TEXTO
    return slide


def add_conteudo_slide(titulo, itens, tabela=None):
    slide = prs.slides.add_slide(LAYOUT_CONTEUDO)
    slide.shapes.title.text = titulo
    slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = COR_TITULO

    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()

    first = True
    for texto, nivel in itens:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = nivel
        _add_rich_text(p, texto)
        for run in p.runs:
            run.font.size = Pt(22 if nivel == 0 else 19)
            run.font.color.rgb = COR_TEXTO

    if tabela:
        _add_tabela(slide, tabela)

    return slide


def _add_rich_text(paragraph, texto):
    partes = texto.split("**")
    for i, parte in enumerate(partes):
        if not parte:
            continue
        run = paragraph.add_run()
        run.text = parte
        run.font.bold = (i % 2 == 1)


def _add_tabela(slide, tabela):
    cabecalho, linhas = tabela
    n_linhas = len(linhas) + 1
    n_colunas = len(cabecalho)

    left = Inches(7.3)
    top = Inches(1.6)
    width = Inches(5.6)
    height = Inches(0.5 * n_linhas)

    shape = slide.shapes.add_table(n_linhas, n_colunas, left, top, width, height)
    table = shape.table

    for j, titulo_col in enumerate(cabecalho):
        cell = table.cell(0, j)
        cell.text = titulo_col
        for p in cell.text_frame.paragraphs:
            p.font.bold = True
            p.font.size = Pt(15)

    for i, linha in enumerate(linhas, start=1):
        for j, valor in enumerate(linha):
            cell = table.cell(i, j)
            _add_rich_text(cell.text_frame.paragraphs[0], valor)
            for run in cell.text_frame.paragraphs[0].runs:
                run.font.size = Pt(15)


def add_secao_final(titulo, subtitulo):
    slide = prs.slides.add_slide(LAYOUT_SECAO)
    slide.shapes.title.text = titulo
    slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(48)
    slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = COR_TITULO
    if len(slide.placeholders) > 1:
        slide.placeholders[1].text = subtitulo
    return slide


# Slide 1 - Capa
add_titulo_slide(
    "Projeto EDA - Ciencia de Dados",
    [
        "Parte A: Percepcao dos Brasileiros Acerca da Democracia",
        "Parte B: Inadimplencia de Pessoas Fisicas no Brasil",
        "",
        "Projeto I - EDA | IMT",
    ],
)

# Slide 2 - Duas perguntas, um fio condutor
add_conteudo_slide(
    "Duas perguntas, um fio condutor",
    [
        ("O que declaramos x o que fazemos", 0),
        ("Parte A: grupos que votam mais tambem **querem** participar mais da politica?", 0),
        ("Parte B: a inadimplencia reage **na hora** aos juros, ou existe um **atraso**?", 0),
        ("Em ambos os casos: **percepcao/decisao declarada vs. comportamento real ao longo do tempo**", 0),
    ],
)

# Slide 3 - Dados e metodo (Parte A)
add_conteudo_slide(
    "Dados e metodo (Parte A)",
    [
        ("Duas bases, uma ponte", 0),
        ("**Ponte:** dimensoes comuns (regiao, escolaridade, idade) -> comparacao ecologica", 1),
        ("**Ferramentas:** Python (pandas, NumPy, pyreadstat), Matplotlib/Seaborn", 1),
        ("**Estatistica:** correlacao de Spearman (rho) entre grupos", 1),
    ],
    tabela=(
        ["Base", "O que mede", "Tamanho"],
        [
            ["CESOP/04832", "Opiniao declarada", "2.000 respondentes"],
            ["TSE 2022", "Comportamento real", "8,8 mi linhas / 311,5 mi aptos"],
        ],
    ),
)

# Slide 4 - O que os brasileiros declaram
add_conteudo_slide(
    "O que os brasileiros declaram",
    [
        ("Quatro retratos da pesquisa (CESOP)", 0),
        ("~65% nao lembram em quem votaram para o Legislativo em 2022", 0),
        ("**Prioridades:** saude (20%), emprego (15%); \"ampliar participacao politica\" e a penultima (2%)", 0),
        ("**Fake news:** preferem **punir (67%)** a **regular (33%)**", 0),
        ("**78% nao tem nenhuma vontade** de participar da politica local", 0),
    ],
)

# Slide 5 - O comportamento real (TSE 2022)
add_conteudo_slide(
    "O comportamento real (TSE 2022)",
    [
        ("Quem comparece? A escolaridade separa", 0),
        ("Comparecimento nacional: **79,4%**", 0),
        ("Por **escolaridade**: de **49% (analfabetos)** a **88% (superior)** -> ~25-38 p.p.", 0),
        ("Por **regiao**: amplitude de **apenas ~3 p.p.**", 0),
    ],
)

# Slide 6 - Nucleo A: escolaridade ALINHA
add_conteudo_slide(
    "Nucleo A: escolaridade ALINHA",
    [
        ("Percepcao e comportamento sobem juntos - rho = 1,00", 0),
        ("A cada degrau de escolaridade sobem juntos:", 0),
        ("Comparecimento real (TSE): ~62% -> ~87%", 1),
        ("Disposicao declarada (CESOP): ~11% -> ~34%", 1),
        ("**Correlacao perfeita (rho = 1,00)** - a renda reproduz o mesmo gradiente", 0),
    ],
)

# Slide 7 - Nucleo A: regiao DISSOCIA
add_conteudo_slide(
    "Nucleo A: regiao DISSOCIA",
    [
        ("Por regiao, percepcao e comportamento se OPOEM - rho = -0,30", 0),
        ("**Sul:** vota mais (~81%), mas e o **mais desinteressado** (~81% \"nenhuma vontade\")", 0),
        ("**Norte:** declara a **maior lembranca de voto (41%)**, mas tem o **menor comparecimento real (78%)**", 0),
        ("**Comparecer != querer participar**", 0),
    ],
)

# Slide 8 - Conclusao Parte A
add_conteudo_slide(
    "Conclusao - Parte A",
    [
        ("E a escolaridade - nao a regiao - que organiza o engajamento", 0),
        ("**Unico eixo** em que opiniao e comportamento se **alinham** (rho = 1,00)", 0),
        ("O brasileiro retratado **comparece** (79,4%), mas **nao quer participar** e **nao lembra** do voto", 0),
        ("-> quer **resultados** mais que **tomar parte**", 1),
    ],
)

# Slide 9 - Dados e metodo (Parte B)
add_conteudo_slide(
    "Dados e metodo (Parte B)",
    [
        ("Serie temporal direto da fonte oficial", 0),
        ("**Fonte:** API do Banco Central (SGS) - series em tempo real", 0),
        ("Inadimplencia PF (21082), Selic (4390), Cambio/Dolar (3695)", 1),
        ("**Periodo:** mar/2011 a abr/2026 - **182 observacoes mensais**", 0),
        ("**Tecnicas:** decomposicao aditiva, correlacao com lags de 3 e 6 meses, teste ADF (estacionariedade), modelo ARIMA(1,1,1)", 0),
    ],
)

# Slide 10 - Evolucao historica e sazonalidade
add_conteudo_slide(
    "Evolucao historica e sazonalidade",
    [
        ("Ciclos macroeconomicos + calendario das familias", 0),
        ("**Ciclos:** alta na crise 2015-2017, queda 2018-2020, **queda atipica na pandemia** (auxilios/renegociacoes), **pico historico em 2026**", 0),
        ("Sazonalidade mensal:", 0),
        ("Pico de estresse: **abril (3,31%) e maio (3,29%)** - \"ressaca\" de IPVA/IPTU/material escolar", 1),
        ("Alivio: **dezembro (3,09%)** - efeito do **13o salario**", 1),
    ],
)

# Slide 11 - Nucleo B: efeito defasado da Selic
add_conteudo_slide(
    "Nucleo B: o efeito defasado da Selic",
    [
        ("A Selic de hoje e o calote de dentro de 6 meses", 0),
        ("**Dolar** tem efeito **fraco e indireto** (via inflacao -> Selic)", 0),
    ],
    tabela=(
        ["Variavel", "Correlacao com inadimplencia"],
        [
            ["Selic (mes corrente)", "0,527"],
            ["Selic - 3 meses", "0,673"],
            ["**Selic - 6 meses**", "**0,774**"],
            ["Dolar (qualquer defasagem)", "-0,21 a -0,28 (fraco)"],
        ],
    ),
)

# Slide 12 - ARIMA: o que vem pela frente
add_conteudo_slide(
    "ARIMA: o que vem pela frente",
    [
        ("Serie nao-estacionaria -> diferenciada -> projecao de alta", 0),
        ("Teste **ADF**: serie original **nao-estacionaria** (p = 0,272); diferenciada **estacionaria** (p = 0,036)", 0),
        ("Modelo **ARIMA(1,1,1)**: coeficientes significativos, residuos sem autocorrelacao (Ljung-Box p = 0,47)", 0),
        ("**Forecast (6 meses): tendencia de alta continua**", 0),
    ],
)

# Slide 13 - Conclusao geral
add_conteudo_slide(
    "Conclusao geral",
    [
        ("Duas defasagens, um Brasil", 0),
        ("**Parte A:** votar e obrigatorio, mas **querer participar** depende da **escolaridade** - nao da regiao", 0),
        ("**Parte B:** o calote de hoje reflete a **Selic de 6 meses atras** - politica monetaria tem efeito **lento**", 0),
        ("**Limitacoes:** comparacao ecologica (A) e modelo univariado (B) - caminhos para trabalhos futuros", 0),
    ],
)

# Slide 14 - Encerramento
add_secao_final("Obrigado!", "Perguntas?")

prs.save("reports/apresentacao.pptx")
print("OK")
